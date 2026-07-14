from __future__ import annotations

import io
import json
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from pypdf import PdfWriter

from content_reader.polish import run_codex_polish
from content_reader.store import VaultStore


def sample_pdf(page_count: int = 2) -> bytes:
    writer = PdfWriter()
    for _ in range(page_count):
        writer.add_blank_page(width=612, height=792)
    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def sample_pptx() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Fourier transform"
    first.placeholders[1].text = "Frequency-domain representation"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Angular frequency"
    second.placeholders[1].text = "omega is measured in radians per second"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


class VaultStoreTest(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        self.vault = self.root / "HKU_Obsidian"
        (self.vault / "Berkeley" / "CS161").mkdir(parents=True)
        self.config = self.root / "config.json"
        self.config.write_text(
            json.dumps(
                {
                    "vault_path": str(self.vault),
                    "notes_root": "Lecture Notes",
                    "route_to_existing_course_folder": True,
                }
            ),
            encoding="utf-8",
        )
        self.store = VaultStore(self.config)
        self.store.ensure_layout()

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def test_pdf_round_trip_and_idempotent_stage_two(self) -> None:
        pdf = sample_pdf()
        record = self.store.import_document(
            filename="memory.pdf",
            content=pdf,
            course="CS161",
            title="Memory Safety",
            lecture_date="2026-07-13",
        )
        self.assertEqual(record["page_count"], 2)
        self.assertTrue(record["source_path"].startswith("Berkeley/CS161/Lecture Notes/"))
        self.assertEqual((self.vault / record["source_path"]).read_bytes(), pdf)

        note = "A canary detects a stack overwrite.\n\nInline math: $\\omega = 2\\pi f$."
        self.store.save_note(record["id"], 2, note)
        self.assertEqual(self.store.get_notes(record["id"])["2"], note)
        pending = self.store.pending_documents()
        self.assertEqual([item["id"] for item in pending], [record["id"]])

        input_hash = self.store.input_hash(record["id"])
        polished_path = self.store.finalize_polished(
            record["id"],
            "## Memory safety\n\nA stack canary detects overwrites before return control is used.",
            expected_hash=input_hash,
        )
        polished = polished_path.read_text(encoding="utf-8")
        raw = (self.vault / record["raw_note_path"]).read_text(encoding="utf-8")
        self.assertIn("Class memos: [[", polished)
        self.assertIn("Polished note: [[", raw)
        self.assertEqual(self.store.pending_documents(), [])

        self.store.save_note(record["id"], 1, "A new memo after Stage 2.")
        self.assertEqual(len(self.store.pending_documents()), 1)

        page = self.store.render_page(record["id"], 1)
        self.assertTrue(page.exists())
        self.assertGreater(page.stat().st_size, 100)

    def test_pptx_extracts_pages_and_uses_central_fallback(self) -> None:
        record = self.store.import_document(
            filename="signals.pptx",
            content=sample_pptx(),
            course="ELEC2441",
            title="Signals",
            lecture_date="2026-07-13",
        )
        self.assertEqual(record["page_count"], 2)
        self.assertTrue(record["source_path"].startswith("Lecture Notes/_Sources/ELEC2441/"))
        extracted = (self.vault / record["extracted_path"]).read_text(encoding="utf-8")
        self.assertIn("Fourier transform", extracted)
        self.assertIn("Angular frequency", extracted)
        page = self.store.render_page(record["id"], 2)
        self.assertTrue(page.exists())

        (self.vault / "ELEC2441").mkdir()
        changes = self.store.reconcile_course_locations()
        self.assertEqual(len(changes), 1)
        moved = self.store.get_document(record["id"])
        self.assertTrue(moved["source_path"].startswith("ELEC2441/Lecture Notes/"))
        self.assertTrue((self.vault / moved["raw_note_path"]).exists())
        raw = (self.vault / moved["raw_note_path"]).read_text(encoding="utf-8")
        self.assertIn(moved["polished_note_path"][:-3], raw)

    def test_stage_two_cannot_write_project_source(self) -> None:
        record = self.store.import_document(
            filename="safe-polish.pdf",
            content=sample_pdf(),
            course="CS161",
            title="Safe Polish",
            lecture_date="2026-07-13",
        )
        self.store.save_note(record["id"], 1, "A memo for the polishing job.")

        captured_command: list[str] = []

        def fake_run(command: list[str], **_kwargs: object) -> SimpleNamespace:
            captured_command.extend(command)
            self.store.finalize_polished(
                record["id"],
                "## Safe result\n\nThe memo is represented.",
                expected_hash=self.store.input_hash(record["id"]),
            )
            return SimpleNamespace(returncode=0)

        def fake_which(executable: str) -> str:
            return "/usr/bin/true" if executable == "codex" else "/usr/bin/python3"

        with patch("content_reader.polish.shutil.which", side_effect=fake_which), patch(
            "content_reader.polish.subprocess.run", side_effect=fake_run
        ):
            result = run_codex_polish(self.store, record["id"])

        self.assertEqual(result["status"], "completed")
        cd_index = captured_command.index("--cd")
        self.assertEqual(Path(captured_command[cd_index + 1]), self.store.drafts_root)
        self.assertIn("--add-dir", captured_command)
        prompt = captured_command[-1]
        self.assertIn("This is a note-polishing task only", prompt)
        self.assertIn("editing code requires the owner's explicit approval", prompt)


if __name__ == "__main__":
    unittest.main()
