from __future__ import annotations

import hashlib
import json
import os
import re
import shutil
import subprocess
import tempfile
import textwrap
import urllib.parse
from datetime import date, datetime
from pathlib import Path
from typing import Any


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_CONFIG_PATH = PROJECT_ROOT / "config.json"


def _now_iso() -> str:
    return datetime.now().astimezone().isoformat(timespec="seconds")


def _safe_segment(value: str, fallback: str) -> str:
    value = re.sub(r"[\\/:*?\"<>|]", "-", value).strip().strip(".")
    value = re.sub(r"\s+", " ", value)
    return value[:120] or fallback


def _json_string(value: str) -> str:
    return json.dumps(value, ensure_ascii=False)


def _atomic_write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, temporary_name = tempfile.mkstemp(
        prefix=f".{path.name}.", suffix=".tmp", dir=path.parent
    )
    try:
        with os.fdopen(fd, "w", encoding="utf-8", newline="\n") as handle:
            handle.write(content)
        os.replace(temporary_name, path)
    finally:
        try:
            os.unlink(temporary_name)
        except FileNotFoundError:
            pass


def _atomic_write_bytes(path: Path, content: bytes) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, temporary_name = tempfile.mkstemp(
        prefix=f".{path.name}.", suffix=".tmp", dir=path.parent
    )
    try:
        with os.fdopen(fd, "wb") as handle:
            handle.write(content)
        os.replace(temporary_name, path)
    finally:
        try:
            os.unlink(temporary_name)
        except FileNotFoundError:
            pass


class StoreError(RuntimeError):
    pass


class VaultStore:
    def __init__(self, config_path: Path | str = DEFAULT_CONFIG_PATH):
        self.config_path = Path(config_path).resolve()
        try:
            self.config = json.loads(self.config_path.read_text(encoding="utf-8"))
        except FileNotFoundError as exc:
            raise StoreError(
                f"Missing configuration file: {self.config_path}. "
                "Copy config.example.json to config.json and choose a notes folder or Obsidian vault."
            ) from exc

        self.storage_mode = str(self.config.get("storage_mode", "obsidian")).strip().lower()
        if self.storage_mode not in {"obsidian", "folder"}:
            raise StoreError("storage_mode must be either 'obsidian' or 'folder'.")
        root_value = (
            self.config.get("vault_path")
            if self.storage_mode == "obsidian"
            else self.config.get("notes_path") or self.config.get("vault_path")
        )
        if not isinstance(root_value, str) or not root_value.strip():
            required_key = "vault_path" if self.storage_mode == "obsidian" else "notes_path"
            raise StoreError(f"Missing required configuration value: {required_key}.")
        self.vault = Path(root_value).expanduser().resolve()
        self.storage_label = "Obsidian vault" if self.storage_mode == "obsidian" else "Markdown folder"
        self.notes_root_name = self.config.get("notes_root", "Lecture Notes")
        self.notes_root = self.vault / self.notes_root_name
        self.runtime_root = PROJECT_ROOT / "runtime"
        self.library_path = self.notes_root / ".content-reader" / "library.json"
        self.extracted_root = self.notes_root / ".content-reader" / "extracted"
        self.rendered_root = self.runtime_root / "rendered"
        self.page_cache_root = self.runtime_root / "pages"
        self.drafts_root = self.runtime_root / "drafts"

    def ensure_layout(self) -> None:
        if self.storage_mode == "folder":
            self.vault.mkdir(parents=True, exist_ok=True)
        elif not self.vault.exists():
            raise StoreError(f"Obsidian vault does not exist: {self.vault}")
        for path in (
            self.notes_root / "_Sources",
            self.notes_root / "Raw",
            self.notes_root / "Polished",
            self.extracted_root,
            self.rendered_root,
            self.page_cache_root,
            self.drafts_root,
        ):
            path.mkdir(parents=True, exist_ok=True)
        if not self.library_path.exists():
            self._save_library({"version": 1, "documents": []})
            self.write_hub([])

    def vault_health(self) -> dict[str, Any]:
        """Prove that the configured notes directory can be read and written."""
        checked_at = _now_iso()
        result: dict[str, Any] = {
            "connected": False,
            "vault_name": self.vault.name,
            "vault_path": str(self.vault),
            "root_name": self.vault.name,
            "root_path": str(self.vault),
            "notes_root": str(self.notes_root),
            "storage_mode": self.storage_mode,
            "storage_label": self.storage_label,
            "obsidian_detected": False,
            "read_verified": False,
            "write_verified": False,
            "checked_at": checked_at,
        }
        try:
            if self.storage_mode == "folder":
                self.vault.mkdir(parents=True, exist_ok=True)
            if not self.vault.is_dir():
                result["reason"] = "The configured notes folder is unavailable."
                return result
            result["obsidian_detected"] = (self.vault / ".obsidian").is_dir()
            if self.storage_mode == "obsidian" and not result["obsidian_detected"]:
                result["reason"] = "The folder does not contain Obsidian configuration."
                return result
            self.notes_root.mkdir(parents=True, exist_ok=True)
            proof_dir = self.notes_root / ".content-reader"
            proof_dir.mkdir(parents=True, exist_ok=True)
            fd, proof_name = tempfile.mkstemp(prefix=".vault-proof-", dir=proof_dir)
            proof_path = Path(proof_name)
            try:
                with os.fdopen(fd, "wb") as handle:
                    handle.write(b"margin-vault-proof")
                result["write_verified"] = True
                result["read_verified"] = proof_path.read_bytes() == b"margin-vault-proof"
            finally:
                proof_path.unlink(missing_ok=True)
            if not result["read_verified"]:
                result["reason"] = "The notes-folder write check could not be read back."
                return result
            result["connected"] = True
            result["reason"] = (
                "Obsidian detected; read and write verified."
                if self.storage_mode == "obsidian"
                else "Markdown folder read and write verified."
            )
            return result
        except OSError as exc:
            result["reason"] = f"Notes-folder verification failed: {exc.strerror or exc}"
            return result

    def _load_library(self) -> dict[str, Any]:
        self.ensure_layout()
        try:
            library = json.loads(self.library_path.read_text(encoding="utf-8"))
        except (FileNotFoundError, json.JSONDecodeError) as exc:
            raise StoreError(f"Cannot read the note library: {exc}") from exc
        if not isinstance(library.get("documents"), list):
            raise StoreError("The note library is malformed.")
        return library

    def _save_library(self, library: dict[str, Any]) -> None:
        _atomic_write_text(
            self.library_path,
            json.dumps(library, ensure_ascii=False, indent=2) + "\n",
        )

    def list_documents(self) -> list[dict[str, Any]]:
        documents = self._load_library()["documents"]
        for document in documents:
            polished = self.vault / document["polished_note_path"]
            document["polished_exists"] = polished.exists()
            document["has_notes"] = self._document_has_notes(document)
            document["polished_current"] = bool(
                polished.exists()
                and document.get("polished_input_hash") == self.input_hash(document)
            )
        return sorted(
            documents,
            key=lambda item: (item.get("lecture_date", ""), item.get("imported_at", "")),
            reverse=True,
        )

    def get_document(self, document_id: str) -> dict[str, Any]:
        for document in self._load_library()["documents"]:
            if document["id"] == document_id:
                return document
        raise StoreError("Lecture not found.")

    def import_document(
        self,
        *,
        filename: str,
        content: bytes,
        course: str,
        title: str,
        lecture_date: str | None = None,
    ) -> dict[str, Any]:
        self.ensure_layout()
        suffix = Path(filename).suffix.lower()
        if suffix not in {".pdf", ".pptx"}:
            raise StoreError("Only PDF and PPTX files are supported.")
        if not content:
            raise StoreError("The selected file is empty.")

        course = _safe_segment(course, "Unsorted")
        title = _safe_segment(title or Path(filename).stem, "Untitled Lecture")
        lecture_date = lecture_date or date.today().isoformat()
        try:
            date.fromisoformat(lecture_date)
        except ValueError as exc:
            raise StoreError("Lecture date must use YYYY-MM-DD.") from exc

        digest = hashlib.sha256(content).hexdigest()
        document_id = hashlib.sha256(
            f"{course}\0{title}\0{digest}".encode("utf-8")
        ).hexdigest()[:16]
        library = self._load_library()
        for existing in library["documents"]:
            if existing["id"] == document_id:
                return existing

        course_destination = self._course_destination(course)
        source_filename = f"{lecture_date} - {title} - {document_id[:6]}{suffix}"
        if course_destination:
            lecture_root = course_destination / "Lecture Notes"
            source_relative = lecture_root / "_Sources" / source_filename
        else:
            lecture_root = Path(self.notes_root_name)
            source_relative = lecture_root / "_Sources" / course / source_filename
        source_path = self.vault / source_relative
        _atomic_write_bytes(source_path, content)

        page_text, page_count = self._extract_source_text(source_path, suffix)
        rendered_pdf = self._prepare_rendered_pdf(source_path, suffix, document_id)

        raw_filename = f"{lecture_date} - {title} - Raw Notes.md"
        polished_filename = f"{lecture_date} - {title} - Polished.md"
        if course_destination:
            raw_relative = lecture_root / "Raw" / raw_filename
            polished_relative = lecture_root / "Polished" / polished_filename
        else:
            raw_relative = lecture_root / "Raw" / course / raw_filename
            polished_relative = lecture_root / "Polished" / course / polished_filename
        extracted_relative = (
            Path(self.notes_root_name)
            / ".content-reader"
            / "extracted"
            / f"{document_id}.md"
        )

        record: dict[str, Any] = {
            "id": document_id,
            "filename": filename,
            "title": title,
            "course": course,
            "course_destination": course_destination.as_posix() if course_destination else None,
            "lecture_date": lecture_date,
            "kind": suffix.lstrip("."),
            "page_count": page_count,
            "source_sha256": digest,
            "source_path": source_relative.as_posix(),
            "raw_note_path": raw_relative.as_posix(),
            "polished_note_path": polished_relative.as_posix(),
            "extracted_path": extracted_relative.as_posix(),
            "rendered_pdf_path": str(rendered_pdf) if rendered_pdf else None,
            "imported_at": _now_iso(),
            "updated_at": _now_iso(),
        }

        _atomic_write_text(self.vault / raw_relative, self._raw_note_template(record))
        _atomic_write_text(
            self.vault / extracted_relative,
            self._extracted_note_template(record, page_text),
        )
        library["documents"].append(record)
        self._save_library(library)
        self.write_hub(library["documents"])
        return record

    def _course_destination(self, course: str) -> Path | None:
        if self.storage_mode != "obsidian":
            return None
        if not self.config.get("route_to_existing_course_folder", True):
            return None
        key = re.sub(r"[^A-Za-z0-9]", "", course).upper()
        if not key:
            return None
        ignored_names = {".obsidian", ".git", ".content-reader", "node_modules"}
        candidates: list[Path] = []
        for root, dirnames, _ in os.walk(self.vault):
            root_path = Path(root)
            dirnames[:] = [
                name
                for name in dirnames
                if name not in ignored_names
                and not name.startswith(".")
                and (root_path / name) != self.notes_root
            ]
            try:
                root_path.relative_to(self.notes_root)
                dirnames[:] = []
                continue
            except ValueError:
                pass
            for name in dirnames:
                if re.sub(r"[^A-Za-z0-9]", "", name).upper() == key:
                    candidates.append((root_path / name).relative_to(self.vault))
        if not candidates:
            return None
        return sorted(candidates, key=lambda path: (len(path.parts), path.as_posix().casefold()))[0]

    def reconcile_course_locations(self) -> list[dict[str, str]]:
        """Re-evaluate course-code routing and repair all affected vault links."""
        library = self._load_library()
        changes: list[dict[str, str]] = []
        for record in library["documents"]:
            destination = self._course_destination(record["course"])
            destination_value = destination.as_posix() if destination else None
            if record.get("course_destination") == destination_value:
                continue

            if destination:
                root = destination / "Lecture Notes"
                new_source = root / "_Sources" / Path(record["source_path"]).name
                new_raw = root / "Raw" / Path(record["raw_note_path"]).name
                new_polished = root / "Polished" / Path(record["polished_note_path"]).name
            else:
                root = Path(self.notes_root_name)
                new_source = root / "_Sources" / record["course"] / Path(record["source_path"]).name
                new_raw = root / "Raw" / record["course"] / Path(record["raw_note_path"]).name
                new_polished = root / "Polished" / record["course"] / Path(record["polished_note_path"]).name

            old_paths = {
                "source_path": record["source_path"],
                "raw_note_path": record["raw_note_path"],
                "polished_note_path": record["polished_note_path"],
            }
            new_paths = {
                "source_path": new_source.as_posix(),
                "raw_note_path": new_raw.as_posix(),
                "polished_note_path": new_polished.as_posix(),
            }
            for key in ("source_path", "raw_note_path", "polished_note_path"):
                old_absolute = self.vault / old_paths[key]
                new_absolute = self.vault / new_paths[key]
                if old_absolute == new_absolute or not old_absolute.exists():
                    continue
                new_absolute.parent.mkdir(parents=True, exist_ok=True)
                if new_absolute.exists():
                    if old_absolute.read_bytes() != new_absolute.read_bytes():
                        raise StoreError(
                            f"Cannot sort {record['title']}: destination already contains a different file named {new_absolute.name}."
                        )
                    old_absolute.unlink()
                else:
                    shutil.move(str(old_absolute), str(new_absolute))

            replacements: list[tuple[str, str]] = []
            for key in ("source_path", "raw_note_path", "polished_note_path"):
                old_value = old_paths[key]
                new_value = new_paths[key]
                replacements.append((old_value, new_value))
                if old_value.lower().endswith(".md"):
                    replacements.append((old_value[:-3], new_value[:-3]))
            replacements.sort(key=lambda pair: len(pair[0]), reverse=True)
            for relative in (new_paths["raw_note_path"], new_paths["polished_note_path"], record["extracted_path"]):
                path = self.vault / relative
                if not path.exists():
                    continue
                text = path.read_text(encoding="utf-8")
                for old_value, new_value in replacements:
                    text = text.replace(old_value, new_value)
                _atomic_write_text(path, text)

            record.update(new_paths)
            record["course_destination"] = destination_value
            if record["kind"] == "pdf":
                record["rendered_pdf_path"] = str(self.vault / new_paths["source_path"])
            changes.append(
                {
                    "document_id": record["id"],
                    "course": record["course"],
                    "from": old_paths["raw_note_path"],
                    "to": new_paths["raw_note_path"],
                }
            )

        if changes:
            self._save_library(library)
            self.write_hub(library["documents"])
        return changes

    def _extract_source_text(
        self, source_path: Path, suffix: str
    ) -> tuple[list[str], int]:
        if suffix == ".pdf":
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(source_path))
                pages = [(page.extract_text() or "").strip() for page in reader.pages]
                return pages, len(reader.pages)
            except Exception as exc:  # PDF parsers surface many vendor exceptions.
                raise StoreError(f"Could not read this PDF: {exc}") from exc

        try:
            from pptx import Presentation

            presentation = Presentation(str(source_path))
            pages: list[str] = []
            for slide in presentation.slides:
                blocks: list[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        blocks.append(text.strip())
                    if getattr(shape, "has_table", False):
                        rows = []
                        for row in shape.table.rows:
                            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                        if rows:
                            blocks.append("\n".join(rows))
                try:
                    notes_frame = slide.notes_slide.notes_text_frame
                    notes = notes_frame.text.strip() if notes_frame else ""
                    if notes:
                        blocks.append(f"Speaker notes:\n{notes}")
                except (AttributeError, ValueError):
                    pass
                pages.append("\n\n".join(dict.fromkeys(blocks)).strip())
            return pages, len(presentation.slides)
        except Exception as exc:
            raise StoreError(f"Could not read this PowerPoint: {exc}") from exc

    def _prepare_rendered_pdf(
        self, source_path: Path, suffix: str, document_id: str
    ) -> Path | None:
        if suffix == ".pdf":
            return source_path

        output_dir = self.rendered_root / document_id
        output_dir.mkdir(parents=True, exist_ok=True)
        final_path = output_dir / "slides.pdf"
        if final_path.exists():
            return final_path

        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            return None
        process = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(source_path),
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
        generated = output_dir / f"{source_path.stem}.pdf"
        if process.returncode != 0 or not generated.exists():
            return None
        generated.replace(final_path)
        return final_path

    def _raw_note_template(self, record: dict[str, Any]) -> str:
        source_label = "Original PDF" if record["kind"] == "pdf" else "Original PowerPoint"
        source_link = self._link(
            record["source_path"], source_label, from_path=record["raw_note_path"]
        )
        polished_link = self._link(
            record["polished_note_path"],
            "Open polished note",
            from_path=record["raw_note_path"],
        )
        lines = [
            "---",
            "content_reader: raw",
            f"document_id: {_json_string(record['id'])}",
            f"course: {_json_string(record['course'])}",
            f"lecture_date: {_json_string(record['lecture_date'])}",
            f"source_sha256: {_json_string(record['source_sha256'])}",
            f"page_count: {record['page_count']}",
            "status: raw",
            f"source_file: {_json_string(record['source_path'])}",
            f"polished_note: {_json_string(record['polished_note_path'])}",
            "tags:",
            "  - lecture-notes/raw",
            "---",
            "",
            f"# {record['title']} — Raw Notes",
            "",
            self._callout_heading("abstract", "Lecture file"),
            f"> Course: **{record['course']}**  ",
            f"> Source: {source_link}  ",
            f"> Polished note: {polished_link}",
            "",
            self._callout_heading("tip", "Page matching"),
            "> Every memo below belongs only to the numbered source page or slide. The marker comments keep the app and Markdown synchronized.",
            "",
        ]
        for page_number in range(1, record["page_count"] + 1):
            lines.extend(
                [
                    f"## Page {page_number}",
                    f"^page-{page_number}",
                    "",
                ]
            )
            if record["kind"] == "pdf":
                page_link = self._link(
                    record["source_path"],
                    f"Open original PDF at page {page_number}",
                    from_path=record["raw_note_path"],
                    page=page_number,
                    embed=self.storage_mode == "obsidian",
                )
                lines.extend(
                    [
                        page_link,
                        "",
                    ]
                )
            elif page_number == 1:
                lines.extend(
                    [
                        source_link,
                        "",
                    ]
                )
            lines.extend(
                [
                    "### Class notes",
                    f"<!-- content-reader:page:{page_number}:start -->",
                    "",
                    f"<!-- content-reader:page:{page_number}:end -->",
                    "",
                ]
            )
        return "\n".join(lines).rstrip() + "\n"

    def _extracted_note_template(
        self, record: dict[str, Any], page_text: list[str]
    ) -> str:
        lines = [
            "---",
            "content_reader: extracted-source",
            f"document_id: {_json_string(record['id'])}",
            f"source_file: {_json_string(record['source_path'])}",
            "---",
            "",
            f"# Extracted source text — {record['title']}",
            "",
            "> Machine-extracted text for the Stage 2 polishing workflow. The original lecture file remains authoritative.",
            "",
        ]
        for index, text in enumerate(page_text, start=1):
            lines.extend([f"## Page {index}", "", text or "_[No extractable text on this page]_", ""])
        return "\n".join(lines).rstrip() + "\n"

    @staticmethod
    def _wiki_target(path: str, keep_suffix: bool = False) -> str:
        value = Path(path).as_posix()
        if not keep_suffix and value.lower().endswith(".md"):
            return value[:-3]
        return value

    def _link(
        self,
        target_path: str,
        label: str,
        *,
        from_path: str,
        page: int | None = None,
        embed: bool = False,
    ) -> str:
        if self.storage_mode == "obsidian":
            target = self._wiki_target(target_path, keep_suffix=not target_path.lower().endswith(".md"))
            fragment = f"#page={page}" if page is not None else ""
            if embed:
                return f"![[{target}{fragment}]]"
            return f"[[{target}{fragment}|{label}]]"

        source_parent = (self.vault / from_path).parent
        target = self.vault / target_path
        relative = Path(os.path.relpath(target, start=source_parent)).as_posix()
        destination = urllib.parse.quote(relative, safe="/#=:@")
        if page is not None:
            destination = f"{destination}#page={page}"
        return f"[{label}]({destination})"

    def _callout_heading(self, kind: str, title: str) -> str:
        if self.storage_mode == "obsidian":
            return f"> [!{kind}] {title}"
        return f"> **{title}**"

    def get_notes(self, document_id: str) -> dict[str, str]:
        record = self.get_document(document_id)
        raw_path = self.vault / record["raw_note_path"]
        text = raw_path.read_text(encoding="utf-8")
        notes: dict[str, str] = {}
        for page_number in range(1, record["page_count"] + 1):
            pattern = re.compile(
                rf"<!-- content-reader:page:{page_number}:start -->\n(.*?)\n<!-- content-reader:page:{page_number}:end -->",
                re.DOTALL,
            )
            match = pattern.search(text)
            notes[str(page_number)] = match.group(1).strip("\n") if match else ""
        return notes

    def save_note(self, document_id: str, page_number: int, content: str) -> dict[str, Any]:
        record = self.get_document(document_id)
        if page_number < 1 or page_number > record["page_count"]:
            raise StoreError("Page number is outside this lecture.")
        content = content.replace("\r\n", "\n").replace("\r", "\n").rstrip()
        if "<!-- content-reader:" in content:
            raise StoreError("That marker is reserved for page synchronization.")

        raw_path = self.vault / record["raw_note_path"]
        text = raw_path.read_text(encoding="utf-8")
        pattern = re.compile(
            rf"(<!-- content-reader:page:{page_number}:start -->)\n.*?\n(<!-- content-reader:page:{page_number}:end -->)",
            re.DOTALL,
        )
        updated, count = pattern.subn(
            lambda match: f"{match.group(1)}\n{content}\n{match.group(2)}",
            text,
            count=1,
        )
        if count != 1:
            raise StoreError("The raw note page marker is missing or duplicated.")
        _atomic_write_text(raw_path, updated)

        library = self._load_library()
        for document in library["documents"]:
            if document["id"] == document_id:
                document["updated_at"] = _now_iso()
                break
        self._save_library(library)
        self.write_hub(library["documents"])
        return {"saved_at": _now_iso(), "has_notes": bool(content.strip())}

    def _document_has_notes(self, record: dict[str, Any]) -> bool:
        try:
            return any(value.strip() for value in self.get_notes(record["id"]).values())
        except (OSError, StoreError):
            return False

    def render_page(self, document_id: str, page_number: int) -> Path:
        return self._render_page_image(document_id, page_number, resolution=120, cache_label="120")

    def render_thumbnail(self, document_id: str, page_number: int) -> Path:
        return self._render_page_image(document_id, page_number, resolution=36, cache_label="thumb-36")

    def _render_page_image(
        self,
        document_id: str,
        page_number: int,
        *,
        resolution: int,
        cache_label: str,
    ) -> Path:
        record = self.get_document(document_id)
        if page_number < 1 or page_number > record["page_count"]:
            raise StoreError("Page number is outside this lecture.")
        cache_dir = self.page_cache_root / document_id
        cache_dir.mkdir(parents=True, exist_ok=True)
        page_path = cache_dir / f"page-{page_number:04d}-{cache_label}.png"
        if page_path.exists():
            return page_path

        pdf_path = record.get("rendered_pdf_path")
        if pdf_path and Path(pdf_path).exists():
            pdftoppm = shutil.which("pdftoppm")
            if not pdftoppm:
                raise StoreError("PDF page renderer is unavailable.")
            prefix = cache_dir / f"page-{page_number:04d}-{cache_label}"
            process = subprocess.run(
                [
                    pdftoppm,
                    "-f",
                    str(page_number),
                    "-l",
                    str(page_number),
                    "-singlefile",
                    "-png",
                    "-r",
                    str(resolution),
                    str(pdf_path),
                    str(prefix),
                ],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if process.returncode == 0 and page_path.exists():
                return page_path

        thumbnail = resolution < 60
        return self._render_text_fallback(
            record,
            page_number,
            page_path,
            size=(480, 270) if thumbnail else (1600, 900),
            thumbnail=thumbnail,
        )

    def _render_text_fallback(
        self,
        record: dict[str, Any],
        page_number: int,
        output_path: Path,
        *,
        size: tuple[int, int] = (1600, 900),
        thumbnail: bool = False,
    ) -> Path:
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError as exc:
            raise StoreError("Slide rendering is unavailable.") from exc

        extracted = (self.vault / record["extracted_path"]).read_text(encoding="utf-8")
        pattern = re.compile(
            rf"^## Page {page_number}\n\n(.*?)(?=\n## Page |\Z)", re.MULTILINE | re.DOTALL
        )
        match = pattern.search(extracted)
        text = (match.group(1).strip() if match else "No extractable slide text.")
        width, height = size
        scale = width / 1600
        image = Image.new("RGB", size, "#f7f4ed")
        draw = ImageDraw.Draw(image)
        font = ImageFont.load_default(size=max(9, round(28 * scale)))
        title_font = ImageFont.load_default(size=max(12, round(40 * scale)))
        draw.rounded_rectangle(
            tuple(round(value * scale) for value in (45, 45, 1555, 855)),
            radius=max(7, round(24 * scale)),
            fill="#ffffff",
            outline="#d9d4c8",
            width=max(1, round(3 * scale)),
        )
        draw.text(
            (round(95 * scale), round(90 * scale)),
            f"{record['title']} · Slide {page_number}",
            fill="#16233a",
            font=title_font,
        )
        wrapped = "\n".join(textwrap.wrap(text, width=88))[:700 if thumbnail else 4200]
        draw.multiline_text(
            (round(95 * scale), round(175 * scale)),
            wrapped,
            fill="#263650",
            font=font,
            spacing=max(4, round(12 * scale)),
        )
        image.save(output_path, format="PNG")
        return output_path

    def pending_documents(self) -> list[dict[str, Any]]:
        pending: list[dict[str, Any]] = []
        for record in self.list_documents():
            if not self._document_has_notes(record):
                continue
            polished_path = self.vault / record["polished_note_path"]
            input_hash = self.input_hash(record)
            if not polished_path.exists() or record.get("polished_input_hash") != input_hash:
                item = dict(record)
                item["input_hash"] = input_hash
                item["raw_note_absolute"] = str(self.vault / record["raw_note_path"])
                item["extracted_absolute"] = str(self.vault / record["extracted_path"])
                item["source_absolute"] = str(self.vault / record["source_path"])
                item["polished_note_absolute"] = str(polished_path)
                pending.append(item)
        return pending

    def input_hash(self, record_or_id: dict[str, Any] | str) -> str:
        record = (
            self.get_document(record_or_id)
            if isinstance(record_or_id, str)
            else record_or_id
        )
        notes = self.get_notes(record["id"])
        payload = json.dumps(notes, ensure_ascii=False, sort_keys=True).encode("utf-8")
        digest = hashlib.sha256()
        digest.update(record["source_sha256"].encode("ascii"))
        digest.update(b"\0")
        digest.update(payload)
        return digest.hexdigest()

    def finalize_polished(
        self, document_id: str, body: str, expected_hash: str | None = None
    ) -> Path:
        record = self.get_document(document_id)
        current_hash = self.input_hash(record)
        if expected_hash and current_hash != expected_hash:
            raise StoreError(
                "Raw notes changed while polishing. The stale draft was not installed; run Stage 2 again."
            )
        body = body.strip()
        if not body:
            raise StoreError("The polished note body is empty.")
        if body.startswith("---"):
            raise StoreError("Provide only the note body; metadata is generated safely.")

        polished_path = self.vault / record["polished_note_path"]
        source_link = self._link(
            record["source_path"],
            "Open lecture file",
            from_path=record["polished_note_path"],
        )
        raw_link = self._link(
            record["raw_note_path"],
            "Open page-linked raw notes",
            from_path=record["polished_note_path"],
        )
        content = "\n".join(
            [
                "---",
                "content_reader: polished",
                f"document_id: {_json_string(record['id'])}",
                f"course: {_json_string(record['course'])}",
                f"lecture_date: {_json_string(record['lecture_date'])}",
                f"source_file: {_json_string(record['source_path'])}",
                f"raw_note: {_json_string(record['raw_note_path'])}",
                "tags:",
                "  - lecture-notes/polished",
                "---",
                "",
                f"# {record['title']}",
                "",
                self._callout_heading("abstract", "Lecture links"),
                f"> Course: **{record['course']}**  ",
                f"> Original source: {source_link}  ",
                f"> Class memos: {raw_link}",
                "",
                body,
                "",
            ]
        )
        _atomic_write_text(polished_path, content)

        raw_path = self.vault / record["raw_note_path"]
        raw = raw_path.read_text(encoding="utf-8")
        raw = raw.replace("status: raw\n", "status: polished-available\n", 1)
        _atomic_write_text(raw_path, raw)

        library = self._load_library()
        for document in library["documents"]:
            if document["id"] == document_id:
                document["polished_input_hash"] = current_hash
                document["polished_at"] = _now_iso()
                break
        self._save_library(library)
        self.write_hub(library["documents"])
        return polished_path

    def write_hub(self, documents: list[dict[str, Any]]) -> Path:
        hub_path = self.notes_root / "Lecture Notes Hub.md"
        hub_relative = hub_path.relative_to(self.vault).as_posix()
        lines = [
            "---",
            "content_reader: hub",
            "tags:",
            "  - lecture-notes",
            "---",
            "",
            "# Lecture Notes Hub",
            "",
            "> Raw class memos remain page-linked to the original lecture file. Polished notes are generated from both the source and those memos.",
            "",
        ]
        courses: dict[str, list[dict[str, Any]]] = {}
        for document in documents:
            courses.setdefault(document["course"], []).append(document)
        if not courses:
            lines.append("_No lectures imported yet._")
        for course in sorted(courses, key=str.casefold):
            lines.extend(
                [
                    f"## {course}",
                    "",
                    "| Date | Lecture | Raw memos | Polished | Source |",
                    "|---|---|---|---|---|",
                ]
            )
            for document in sorted(
                courses[course], key=lambda item: item["lecture_date"], reverse=True
            ):
                polished_exists = (self.vault / document["polished_note_path"]).exists()
                polished_label = "Open" if polished_exists else "Pending"
                raw_link = self._link(document["raw_note_path"], "Raw", from_path=hub_relative)
                polished_link = self._link(
                    document["polished_note_path"], polished_label, from_path=hub_relative
                )
                source_link = self._link(
                    document["source_path"], document["kind"].upper(), from_path=hub_relative
                )
                lines.append(
                    f"| {document['lecture_date']} | {document['title']} | {raw_link} | {polished_link} | {source_link} |"
                )
            lines.append("")
        _atomic_write_text(hub_path, "\n".join(lines).rstrip() + "\n")
        return hub_path
